
import tkinter as tk
from tkinter import ttk
from pptx import Presentation
from tkinter import messagebox
from pptx.util import Pt  # Import Pt to set font size
from pptx.dml.color import RGBColor

import pandas as pd
import os
from datetime import datetime
import re



current_select_bible_name= "개역개정"

말씀 = '창'
장 = [1, 2]
절 = [9]

df = pd.read_csv(f"./bible/{current_select_bible_name}.csv")
df['절'] = pd.to_numeric(df['절'])

condition = (df.색인 == 말씀) & (df.장 >= 장[0])  & (df.장 <= 장[1])
filtered_table = df.loc[ condition ,['색인','장','절','내용']]

print(filtered_table)


## 현재 만들어진 슬라이드 레이아웃에 접근할 수 있는 아이템 index 보는법

from pptx import Presentation # 라이브러리 
from pptx.util import Inches # 사진, 표등을 그리기 위해

prs = Presentation('template.pptx')
for i in range(0, 11):
    print("--------[%d] ------ "%(i))
    slide = prs.slides.add_slide(prs.slide_layouts[i])
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))