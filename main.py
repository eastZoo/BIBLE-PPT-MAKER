import tkinter as tk
from tkinter import ttk
from pptx import Presentation
from tkinter import messagebox
from pptx.util import Pt  # Import Pt to set font size
from pptx.dml.color import RGBColor

import pandas as pd
import os
from datetime import datetime
import re

#custom modules
import query
import pages.history as history

# 초기 쿼리 세팅
query.initial()

rootdir ="/"

root = tk.Tk()
root.title('말씀 ppt 생성기 - made by eastzoo')
root.geometry("510x360")
root.resizable(0, 0)

bible_range_text = tk.StringVar()

#STYLE
basicFont = 10
basicWidth = 18

# 셀렉트 박스 리스트 담는 변수
bible_list_name = []
# 현재 셀렉트 박스에서 선택된 이름 담는 변수
current_select_bible_name =""
# 사용자가 입력한 성경 범위 저장 함수
bible_range = ""

  # 입력받은 문자를 말씀, 장, 절로 분리하는 함수
def extract_numbers(input_string):
    # 첫번째로 입력받은 문자열의 공백을 모두 제거
    input_string = input_string.replace(" ", "")

    # 숫자가 나오기 전까지의 문자열 추출하여 말씀에 저장
    말씀 = ''
    장 = []
    절 = []
    

    # 정규표현식을 사용하여 숫자와 숫자 이후에 ":" 뒤에 나오는 숫자 추출
    pattern = re.compile(r'([^\d:]*)(\d+(?:-\d+)?)(?:[:](\d+(?:-\d+)?))?')
    matches = pattern.match(input_string)

    if matches:
        말씀 = matches.group(1).strip()
        nums = matches.group(2).split('-')
        장.extend([int(num) for num in nums])

        if matches.group(3):
            nums = matches.group(3).split('-')
            절.extend([int(num) for num in nums])
    else:
        말씀 = input_string
    return 말씀, 장, 절

# 입력한 말씀 조건에 맞는 파워포인트 생성 함수
def create_ppt():
    prs = Presentation('template.pptx') # 파워포인트를 다루는 함수 변수 선언Load the PowerPoint template   
   
    # 선택된 성경 .csv 파일 읽어오기
    df = pd.read_csv(f"./bible/{current_select_bible_name}.csv")
    df['절'] = pd.to_numeric(df['절'])

    
    # 빈 DataFrame 생성
    main_df = pd.DataFrame(columns=['색인', '장','절','내용'])

    ### 입력한 성경 필터 처리 가공 시작
    
    result_list = bible_range_text.get().split(',')
    

    # 입력받은 문자를 말씀, 장, 절로 분리하는 함수
    # , 로 받은 복수개의 말씀들을 범위 분해
    for input_bible in result_list:
        말씀, 장, 절 = extract_numbers(input_bible)
        
        # 예) 
        #     창1       =  창세기 1장 전체
        #     롬1-3    =  로마서 1장 1절 - 3장 전체
        #     전1:3    =  전도서 1장 3절
        #     레1-3:9 =  레위기 1장 1절 - 3장 9절
        #     스1:3-9 =  에스라 1장 3절 - 1장 9절

        # ex)  창1
        if(len(장)==1 and len(절)==0):
            print("창1")
            condition = (df.색인 == 말씀) & (df.장 == 장[0])
        # ex) 롬1-3 o
        elif(len(장) > 1 and len(절)==0):
            print("롬1-3")
            condition = (df.색인 == 말씀) & (df.장 >= 장[0])  & (df.장 <= 장[1])
        # # ex) 전1:3 o
        elif(len(장)==1 and len(절)==1):
            print("전1:3")
            condition = (df.색인 == 말씀) & (df.장 == 장[0])  & (df.절 == 절[0])
        # # ex) 레1-3:9
        elif(len(장) > 1 and len(절)==1):
            print("레1-3:9")
            condition_1 = (df['색인'] == 말씀) & (df['장'] >= 장[0])& (df['장'] < 장[1])
            condition_2 = (df['색인'] == 말씀) & (df['장'] == 장[1]) & (df['절'] <= 절[0])
            
            condition = condition_1 | condition_2
        # # ex) 스1:3-9
        elif(len(장)==1 and len(절)>1):
            print("스1:3-9")    
            condition = (df.색인 == 말씀) & (df.장 == 장[0]) & (df.절 >= 절[0]) & (df.절 <= 절[1])
        else:
            messagebox.showinfo("Fail", "잘못된 입력 입니다.")
 
        
        # filtered_table = df.loc[ condition ,['색인','장','절','내용']]
        main_df = pd.concat([main_df, df.loc[ condition ,['색인','장','절','내용']]])

    # 첫번째 페이지
    for index, row in main_df.iterrows():
      # 슬라이드 추가
        slide_layout = prs.slide_layouts[0]  # 0은 제목과 내용이 있는 슬라이드 레이아웃
        slide = prs.slides.add_slide(slide_layout)
        
        # 제목과 내용 추가
        title = slide.shapes.title
        content = slide.placeholders[10]
        
        title.text = f"{row['색인']} {row['장']}장 {row['절']}절"
        
         # 글꼴 크기 및 글꼴 패밀리 설정
        # font = title.font
        # font.name = '나눔스퀘어 ExtraBold'
        # font.size = Pt(31)
        # # 글자색 설정
        # font.color.rgb = RGBColor(255,255,255)
        
        content.text = row['내용']
                
                
    # 현재 날짜와 시간을 형식화하여 가져옵니다.
    today_datetime = datetime.today().strftime("%y%m%d_%H%m%S")
    # Modify the file name to today's date
    name =bible_range_text.get().replace(" ", "")
    new_file_name = f"[{today_datetime}]말씀.pptx"
    
    download_folder = os.path.join(os.path.expanduser('~'), 'Downloads')
    
    # exportPPT 폴더 경로
    export_folder = os.path.join(download_folder, 'exportPPT')

    # exportPPT 폴더가 없으면 생성
    if not os.path.exists(export_folder):
        os.makedirs(export_folder)
        
    # backup 폴더 경로
    backup_folder = "./backup"

    # backup 폴더가 없으면 생성
    if not os.path.exists(backup_folder):
        os.makedirs(backup_folder)


    # 파일 경로 설정
    file_path = os.path.join(export_folder, new_file_name)
    file_path_backup = os.path.join('./backup', new_file_name)
  
    print(file_path)
    # PPT 저장
    prs.save(file_path)
    prs.save(file_path_backup)
    
    query.insert(bible_range_text.get().replace(" ", ""),file_path_backup )
    messagebox.showinfo("Success", "말씀 ppt를 성공적으로 생성했습니다.!")


def delete_page():
    for frame in main_frame.winfo_children():
        frame.destroy()
        
        
# 셀렉트 박스 데이터 가져오는 함수
def get_bible_select_data():
    global bible_list_name
    # 첫시작 시 초기화 코드 ( 화면 실행 될때마다 데이터 쌓임 방지 )
    
    bible_list_name = []
    # 성경 폴더에 있는 성경종류 리스트 가져오기 ( 개역개정, + a 예정 )
    bible_list = os.listdir('./bible')
    # 엑셀 확장자 삭제후 보관할 리스트 변수
    for file in bible_list:
        if file.count(".") == 1: 
            name = file.split('.')[0]
            bible_list_name.append(name)
        else:
            for k in range(len(file)-1,0,-1):
                if file[k]=='.':
                    bible_list_name.append(file[:k])
                    break
    return bible_list_name
    
    
# home 페이지 띄우는 함수
def home_page():
    #셀렉트 박스 정보 가져 오기
    get_bible_select_data()
    
    global bible_list_name
    global current_select_bible_name
    
    # 현재 선택한 성경의 변수 저장
    current_select_bible_name = bible_list_name[0]
  
    home_frame = tk.Frame(main_frame)
                    
    # start 셀렉트박스 wrapper 
    ## 성경 셀렉트박스 라벨
    bibile_label = tk.Label(home_frame, text='성경', font=('Bold', basicFont))
    bibile_label.place(x=10, y= 1)

    ## 성경 셀렉트박스
    select_bible = ttk.Combobox(home_frame, values=bible_list_name, width=basicWidth-2)
    select_bible.set(bible_list_name[0])
    select_bible.place(x=80, y=1)
    
    
     # 셀렉트 박스에서 현재 선택된 값 변수에 저장 
    def getSelectedItem(arg):
        global current_select_bible_name
        # 선택한 항목 테이블 만드는 함수
        current_select_bible_name = select_bible.get()
    select_bible.bind("<<ComboboxSelected>>", getSelectedItem)
    
    # end 셀렉트박스 wrapper 
    
    
    # start 성경 범위 입력
    
    def on_enter_press(event):
        create_ppt()
        
    ## 성경범위 라벨
    bibile_range = tk.Label(home_frame, text='성경 구절', font=('Bold', basicFont))
    bibile_range.place(x=10, y= 30)
    
    example_text = tk.Label(home_frame, text='예) 창1:1-3 , 스1:3-5, 스1 ( "," 콤마구분으로 말씀 ppt 동시 생성가능 )', font=('Bold', basicFont-2))
    example_text.place(x=10, y= 57)
    
    ## 성경범위 input
    bibile_range_textbox = ttk.Entry(home_frame, textvariable=bible_range_text,font=('Bold', basicFont), width=basicWidth+1)
    # 엔터 키를 눌렀을 때 on_enter_press 함수 실행
    bibile_range_textbox.bind("<Return>", on_enter_press)
    bibile_range_textbox.place(x=80, y= 30)
    bibile_range_textbox.focus()
    # end 성경 범위 입력
    
    
    
    ## 성경 인덱스 리스트
    
    # 리스트 박스 생성
    # padx = 좌우 패딩
    bible_index_list = query.bible_index_select()
    
    bible_index_tree = ttk.Treeview(home_frame, column=("c1", "c2"), show='headings', height=5)
    bible_index_tree.place(x=10, y= 80)
    
    bible_index_tree.column("# 1", anchor="center",width=150)
    bible_index_tree.heading("# 1", text="성경")
    bible_index_tree.column("# 2", anchor="center",width=150)
    bible_index_tree.heading("# 2", text="약어")
    
    
    def on_click(event):
        print("selected items:")
        item = bible_index_tree.focus()

        current_text = bibile_range_textbox.get()
        if(len(current_text) == 0):
            bibile_range_textbox.insert(tk.END, current_text + f'{item}')  
            return
        
        bibile_range_textbox.delete(0, tk.END)  # Entry 초기화
        bibile_range_textbox.insert(tk.END, current_text + f', {item}')  
        print(item)
  
        
              
    vsb = ttk.Scrollbar(home_frame, orient="vertical", command=bible_index_tree.yview)
    vsb.place(x=294, y= 81, height=122)
    bible_index_tree.configure(yscrollcommand=vsb.set)
    
    # 데이터를 리스트 박스에 추가
    for item in bible_index_list:
        bible_index_tree.insert('',"end", values=(f'{item[0]}',f'{item[1]}'), iid=item[2])
    
    # 트리뷰에 이벤트 바인딩
    bible_index_tree.bind("<<TreeviewSelect>>", on_click)
    
    
    
    Fact="""예)
창1      =  창세기 1장 전체
롬1-3    =  로마서 1장 1절 - 3장 전체
전1:3    =  전도서 1장 3절
레1-3:9  =  레위기 1장 1절 - 3장 9절
스1:3-9  =  에스라 1장 3절 - 1장 9절
    """
    # 사1:3-3:9= 이사야 1장 3절 - 3장 9절
    
    T = tk.Text(home_frame, height = 8, width =42, )
    T.place(x=10, y=220)
    T.insert(tk.END, Fact)
    
    make_ppt_btn = tk.Button(home_frame, text='PPT 만들기', font=('Bold', basicFont), height=3, command=create_ppt)
    make_ppt_btn.place(x=230, y=1)

    home_frame.pack(expand=True, fill='both', pady=20)



# 사이드 메뉴 선택에 따라 다른 페이지 숨김 함수
def hide_indicators():
    home_indicate.config(bg='#DADADA')
    history_indicate.config(bg='#DADADA')
    
# 페이지 라우팅미들웨어 + 사이드 메뉴 클릭시 클릭 css 표시 함수
def indicate(type,lb, page):
    hide_indicators()
    lb.config(bg='#F15642')
    delete_page()
    
    # 페이지 별 라우팅
    if(type == "history"):
        print("HISTORY")
        print(page)
        page()
    if(type == "home"):
        page()

# 사이드바  options_frame
options_frame = tk.Frame(root, bg='#DADADA')

options_frame.pack(side=tk.LEFT)
options_frame.pack_propagate(False)
options_frame.configure(width=50, height=400)

# 메인 프레임  main_frame
main_frame = tk.Frame(root, highlightthickness=2)

main_frame.pack(side=tk.LEFT)
main_frame.pack_propagate(False)
main_frame.configure(height=400, width=500)


# 홈 버튼 이미지
#start
homeImg = tk.PhotoImage(file="./assets/icons/ppt.png")
home_btn = tk.Button(options_frame, text='Home', font=('Bold', 15), image=homeImg, height=50, width=50, command=lambda: indicate("home",home_indicate, home_page))
home_btn.place(x=0, y=0)

home_indicate = tk.Label(options_frame, text='')
home_indicate.place(x=0, y=0, width=5, height=55)
#end

    
    
# 기록(history) 버튼 이미지
#start
historyImg = tk.PhotoImage(file="./assets/icons/history.png")
history_btn = tk.Button(options_frame, text='history', font=('Bold', 15), 
                        image=historyImg, height=50, width=50, command=lambda: indicate("history",history_indicate, lambda: history.history_page(main_frame)))
history_btn.place(x=0, y=55)

history_indicate = tk.Label(options_frame, text='')
history_indicate.place(x=0, y=55, width=5, height=55)
#end


# 처음 시작할때 메인 페이지 세팅 __INIT__
indicate("home",home_indicate, home_page)
root.mainloop() 