import tkinter as tk
from tkinter import ttk
from pptx import Presentation
from tkinter import messagebox
from pptx.util import Pt  # Import Pt to set font size
from pptx.dml.color import RGBColor

import pandas as pd
import os
from datetime import datetime
import re
#custom modules
import query

# prev 페이지 띄우는 함수
def history_page(main_frame):
    history_list = query.select()
    history_frame = tk.Frame(main_frame)
    
     # 입력받은 문자를 말씀, 장, 절로 분리하는 함수
    def extract_numbers(input_string):
        # 첫번째로 입력받은 문자열의 공백을 모두 제거
        input_string = input_string.replace(" ", "")

        # 숫자가 나오기 전까지의 문자열 추출하여 말씀에 저장
        말씀 = ''
        장 = []
        절 = []
        

        # 정규표현식을 사용하여 숫자와 숫자 이후에 ":" 뒤에 나오는 숫자 추출
        pattern = re.compile(r'([^\d:]*)(\d+(?:-\d+)?)(?:[:](\d+(?:-\d+)?))?')
        matches = pattern.match(input_string)

        if matches:
            말씀 = matches.group(1).strip()
            nums = matches.group(2).split('-')
            장.extend([int(num) for num in nums])

            if matches.group(3):
                nums = matches.group(3).split('-')
                절.extend([int(num) for num in nums])
        else:
            말씀 = input_string
        return 말씀, 장, 절



    # 입력한 말씀 조건에 맞는 파워포인트 생성 함수
    def create_ppt(data):
        prs = Presentation('template.pptx') # 파워포인트를 다루는 함수 변수 선언Load the PowerPoint template  
        
        #일단 개역개정
        current_select_bible_name = "개역개정" 
    
        # 선택된 성경 .csv 파일 읽어오기
        df = pd.read_csv(f"./bible/{current_select_bible_name}.csv")
        df['절'] = pd.to_numeric(df['절'])

        
        # 빈 DataFrame 생성
        main_df = pd.DataFrame(columns=['색인', '장','절','내용'])

        ### 입력한 성경 필터 처리 가공 시작
        
        result_list = data.split(',')
        

        # 입력받은 문자를 말씀, 장, 절로 분리하는 함수
        # , 로 받은 복수개의 말씀들을 범위 분해
        for input_bible in result_list:
            말씀, 장, 절 = extract_numbers(input_bible)
            
            # 예) 
            #     창1       =  창세기 1장 전체
            #     롬1-3    =  로마서 1장 1절 - 3장 전체
            #     전1:3    =  전도서 1장 3절
            #     레1-3:9 =  레위기 1장 1절 - 3장 9절
            #     스1:3-9 =  에스라 1장 3절 - 1장 9절

            # ex)  창1
            if(len(장)==1 and len(절)==0):
                print("창1")
                condition = (df.색인 == 말씀) & (df.장 == 장[0])
            # ex) 롬1-3 o
            elif(len(장) > 1 and len(절)==0):
                print("롬1-3")
                condition = (df.색인 == 말씀) & (df.장 >= 장[0])  & (df.장 <= 장[1])
            # # ex) 전1:3 o
            elif(len(장)==1 and len(절)==1):
                print("전1:3")
                condition = (df.색인 == 말씀) & (df.장 == 장[0])  & (df.절 == 절[0])
            # # ex) 레1-3:9
            elif(len(장) > 1 and len(절)==1):
                print("레1-3:9")
                condition_1 = (df['색인'] == 말씀) & (df['장'] >= 장[0])& (df['장'] < 장[1])
                condition_2 = (df['색인'] == 말씀) & (df['장'] == 장[1]) & (df['절'] <= 절[0])
                
                condition = condition_1 | condition_2
            # # ex) 스1:3-9
            elif(len(장)==1 and len(절)>1):
                print("스1:3-9")    
                condition = (df.색인 == 말씀) & (df.장 == 장[0]) & (df.절 >= 절[0]) & (df.절 <= 절[1])
            else:
                messagebox.showinfo("Fail", "잘못된 입력 입니다.")
    
            
            # filtered_table = df.loc[ condition ,['색인','장','절','내용']]
            main_df = pd.concat([main_df, df.loc[ condition ,['색인','장','절','내용']]])

        # 첫번째 페이지
        for index, row in main_df.iterrows():
        # 슬라이드 추가
            slide_layout = prs.slide_layouts[0]  # 0은 제목과 내용이 있는 슬라이드 레이아웃
            slide = prs.slides.add_slide(slide_layout)
            
            # 제목과 내용 추가
            title = slide.shapes.title
            content = slide.placeholders[10]
            
            title.text = f"{row['색인']} {row['장']}장 {row['절']}절"
            
            # 글꼴 크기 및 글꼴 패밀리 설정
            # font = title.font
            # font.name = '나눔스퀘어 ExtraBold'
            # font.size = Pt(31)
            # # 글자색 설정
            # font.color.rgb = RGBColor(255,255,255)
            
            content.text = row['내용']
                    
                    
        # 현재 날짜와 시간을 형식화하여 가져옵니다.
        today_datetime = datetime.today().strftime("%y%m%d_%H%m%S")
        # Modify the file name to today's date
        name =data.replace(" ", "")
        new_file_name = f"[{today_datetime}]말씀.pptx"
        
        download_folder = os.path.join(os.path.expanduser('~'), 'Downloads')
        
        # exportPPT 폴더 경로
        export_folder = os.path.join(download_folder, 'exportPPT')
        # exportPPT 폴더가 없으면 생성
        if not os.path.exists(export_folder):
            os.makedirs(export_folder)

        # 파일 경로 설정
        file_path = os.path.join(export_folder, new_file_name)

        # PPT 저장
        prs.save(file_path)
        messagebox.showinfo("Success", "말씀 ppt를 성공적으로 생성했습니다.!")


    
    def show_selected_item():
        # 선택한 아이템의 ID 가져오기
        selected_item = ListTree.selection()
        if selected_item:
            # 선택한 아이템의 정보 가져오기
            item = ListTree.item(selected_item)
            values = item['values']
            print("Selected item:", values)
            
            create_ppt(values[1])
        else:
            print("No item selected")
    
    # 리스트 리패칭
    def refresh_list():
        # 현재 트리뷰의 모든 항목 삭제
        ListTree.delete(*ListTree.get_children())
        # 삭제 후 리패칭
        history_list = query.select()
     
        # 데이터를 리스트 박스에 추가
        for item in history_list:
            ListTree.insert('',"end", values=(f'{item[0]}',f'{item[1]}',f'{item[2]}'))

    def delete_selected_item():
        selected_item = ListTree.selection()
        if selected_item:
            # 선택한 아이템의 정보 가져오기
            item = ListTree.item(selected_item)
            values = item['values']
            # id를 uniq 키로 전송하여 삭제
            query.delete(values[0])
        else:
            print("No item selected")
        #삭제 리패칭
        refresh_list()
    
    # 리스트 박스 생성
    # padx = 좌우 패딩
    ListTree = ttk.Treeview(history_frame, column=("c1", "c2", "c3"), show='headings', height=5)
    ListTree.pack(side="left", fill="y",padx=10)
    
    ListTree.column("# 1", anchor="center", width=30)
    ListTree.heading("# 1", text="id")
    ListTree.column("# 2", anchor="center",width=150)
    ListTree.heading("# 2", text="내용")
    ListTree.column("# 3", anchor="center",width=150)
    ListTree.heading("# 3", text="날짜")

    # 데이터를 리스트 박스에 추가
    for item in history_list:
        ListTree.insert('',"end", values=(f'{item[0]}',f'{item[1]}',f'{item[3]}',f'{item[2]}'))

    # 확인 버튼 생성
    button = tk.Button(history_frame, text="불러오기(저장)", command=show_selected_item)
    button.pack(side="top", fill="x")
    
     # 삭제 버튼 생성
    button = tk.Button(history_frame, text="삭제", command=delete_selected_item)
    button.pack(side="top", fill="x")

    history_frame.pack(expand=True, fill='both', pady=20)

    